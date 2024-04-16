from pptx import Presentation
from pptx.util import Inches


def insert_image_into_ppt_slide(ppt_path, slide_index, image_path, position, width, height, margin=0.5):
    """
    Insert an image into a specified position on a specified slide of a PowerPoint file with a margin.

    Args:
    - param ppt_path: str - Path to the PowerPoint file.
    - param slide_index: int - Index of the slide (0-based).
    - param image_path: str - Path to the image file.
    - param position: str - Position for the image ('top_left', 'top_center', 'top_right',
    -                 'middle_left', 'middle_right', 'bottom_left', 'bottom_center', 'bottom_right').
    - param width: float - Width of the image when inserted (in inches).
    - param height: float - Height of the image when inserted (in inches).
    - param margin: float - Margin around the image (in inches), default is 0.5 inch.
    Returns:
    - None: The function directly modifies the pptx file.

    
    """
        # Open the PowerPoint file
    try:
        prs = Presentation(ppt_path)
    except IOError:
        raise FileNotFoundError("The specified PowerPoint file does not exist.")

    # Check if the slide index is in range
    if slide_index >= len(prs.slides):
        raise IndexError("Slide index out of range.")

    # Get the specified slide
    slide = prs.slides[slide_index]

    # Get slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Calculate position for the image with margin
    margin_inch = Inches(margin)
    positions = {
        'top_left': (margin_inch, margin_inch),
        'top_center': (slide_width / 2 - Inches(width) / 2, margin_inch),
        'top_right': (slide_width - Inches(width) - margin_inch, margin_inch),
        'middle_left': (margin_inch, slide_height / 2 - Inches(height) / 2),
        'middle_right': (slide_width - Inches(width) - margin_inch, slide_height / 2 - Inches(height) / 2),
        'bottom_left': (margin_inch, slide_height - Inches(height) - margin_inch),
        'bottom_center': (slide_width / 2 - Inches(width) / 2, slide_height - Inches(height) - margin_inch),
        'bottom_right': (slide_width - Inches(width) - margin_inch, slide_height - Inches(height) - margin_inch)
    }

    if position not in positions:
        raise ValueError("Invalid position. Choose from 'top_left', 'top_center', 'top_right', 'middle_left', "
                        "'middle_right', 'bottom_left', 'bottom_center', 'bottom_right'.")

    left, top = positions[position]

    # Insert the image
    try:
        slide.shapes.add_picture(image_path, left, top, width=Inches(width), height=Inches(height))
    except FileNotFoundError:
        raise FileNotFoundError("The specified image file does not exist.")

    # Save the changes
    prs.save(ppt_path)
