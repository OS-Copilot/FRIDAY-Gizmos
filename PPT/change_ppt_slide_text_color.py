from pptx import Presentation
from pptx.dml.color import RGBColor


def change_ppt_slide_text_color(ppt_path, slide_index, text_type, color):
    """
    Change the color of the title or body text in a specified slide of a PowerPoint file.

    Args:
    - param ppt_path: str - Path to the PowerPoint file.
    - param slide_index: int - Index of the slide (0-based).
    - param text_type: str - Type of text to change ('title' or 'body').
    - param color: tuple - New color as a tuple (R, G, B).
    
    Returns:
    - None: The function directly modifies the pptx file.

    """
    # Open the PowerPoint file
    try:
        prs = Presentation(ppt_path)
    except IOError:
        raise FileNotFoundError("The specified PowerPoint file does not exist.")

    # Check if the slide index is in range
    if slide_index >= len(prs.slides):
        raise IndexError("Slide index out of range.")

    # Get the specified slide
    slide = prs.slides[slide_index]

    # Function to change text color in a text frame
    def change_text_frame_color(text_frame):
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(*color)

    # Change the color of the specified text type
    if text_type == 'title':
        if slide.shapes.title:
            change_text_frame_color(slide.shapes.title.text_frame)
        else:
            raise ValueError("The specified slide does not have a title.")
    elif text_type == 'body':
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape != slide.shapes.title:
                change_text_frame_color(shape.text_frame)
    else:
        raise ValueError("Invalid text type. Please use 'title' or 'body'.")

    # Save the changes
    prs.save(ppt_path)
