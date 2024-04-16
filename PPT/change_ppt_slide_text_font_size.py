from pptx import Presentation
from pptx.util import Pt


def change_ppt_slide_text_font_size(ppt_path, slide_index, text_type, font_size):
    """
    Change the font size of the title or body text in a specified slide of a PowerPoint file.

    Args:
    - param ppt_path: str - Path to the PowerPoint file.
    - param slide_index: int - Index of the slide (0-based).
    - param text_type: str - Type of text to change ('title' or 'body').
    - param font_size: int - New font size.
    
    Returns:
    - None: The function directly modifies the pptx file.

    """
    # Open the PowerPoint file
    try:
        prs = Presentation(ppt_path)
    except IOError:
        raise FileNotFoundError("The specified PowerPoint file does not exist.")

    # Check if the slide index is in range
    if slide_index >= len(prs.slides):
        raise IndexError("Slide index out of range.")

    # Get the specified slide
    slide = prs.slides[slide_index]

    # Function to change font size in a text frame
    def change_text_frame_font_size(text_frame, font_size):
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size)

    # Change the font size of the specified text type
    if text_type == 'title':
        if slide.shapes.title:
            change_text_frame_font_size(slide.shapes.title.text_frame, font_size)
        else:
            raise ValueError("The specified slide does not have a title.")
    elif text_type == 'body':
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape != slide.shapes.title:
                change_text_frame_font_size(shape.text_frame, font_size)
    else:
        raise ValueError("Invalid text type. Please use 'title' or 'body'.")

    # Save the changes
    prs.save(ppt_path)
