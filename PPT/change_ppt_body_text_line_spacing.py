from pptx import Presentation


def change_ppt_body_text_line_spacing(ppt_path, slide_index, line_spacing):
    """
    Change the line spacing of the body text in a specified slide of a PowerPoint file.

    Args:
    - param ppt_path: str - Path to the PowerPoint file.
    - param slide_index: int - Index of the slide (0-based).
    - param line_spacing:  Line spacing value. Use a float for multiple spacing (e.g., 1.5 for 1.5x spacing),
                        or Pt for exact spacing in points.
    
    Returns:
    - None: The function directly modifies the pptx file.
    """
    # Open the PowerPoint file
    try:
        prs = Presentation(ppt_path)
    except IOError:
        raise FileNotFoundError("The specified PowerPoint file does not exist.")

    # Check if the slide index is in range
    if slide_index >= len(prs.slides):
        raise IndexError("Slide index out of range.")

    # Get the specified slide
    slide = prs.slides[slide_index]

    # Change line spacing of each body text shape
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            if isinstance(line_spacing, float):
                # Set multiple line spacing (e.g., 1.5 lines)
                paragraph.line_spacing = line_spacing
            else:
                # Set exact line spacing
                paragraph.line_spacing = line_spacing

    # Save the changes
    prs.save(ppt_path)
